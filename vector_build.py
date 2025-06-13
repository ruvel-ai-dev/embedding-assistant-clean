import os
import re
import io
import hashlib
import json
from azure.storage.blob import BlobServiceClient
from concurrent.futures import ThreadPoolExecutor, as_completed
from docx import Document as DocxDocument
from pptx import Presentation
from PyPDF2 import PdfReader
from langchain_community.vectorstores import FAISS
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
from dotenv import load_dotenv

load_dotenv()

AZURE_CONNECTION_STRING = os.getenv("AZURE_STORAGE_CONNECTION_STRING")
AZURE_CONTAINER_NAME = "resources"
INDEX_FILE = "faiss_index"
HASH_RECORD_FILE = "hashes.json"
MAX_WORKERS = int(os.getenv("CONCURRENT_WORKERS", "3"))

# Load existing hash records if present
if os.path.exists(HASH_RECORD_FILE):
    with open(HASH_RECORD_FILE, "r") as f:
        existing_hashes = json.load(f)
else:
    existing_hashes = {}


def calculate_file_hash(content):
    return hashlib.md5(content).hexdigest()


def get_blob_content(blob_client):
    stream = io.BytesIO()
    blob_client.download_blob().readinto(stream)
    stream.seek(0)
    return stream.read()


def read_file(file_bytes, filename):
    if filename.endswith(".pdf"):
        reader = PdfReader(io.BytesIO(file_bytes))
        return "\n".join([page.extract_text() or "" for page in reader.pages])
    elif filename.endswith(".docx"):
        doc = DocxDocument(io.BytesIO(file_bytes))
        return "\n".join([para.text for para in doc.paragraphs])
    elif filename.endswith(".pptx"):
        prs = Presentation(io.BytesIO(file_bytes))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    elif filename.endswith(".txt"):
        return file_bytes.decode("utf-8", errors="ignore")
    else:
        raise ValueError("Unsupported file type")


def is_general_purpose(summary, tags, filename):
    name_lower = filename.lower()
    return ("general" in name_lower or "main" in name_lower
            or "general" in summary.lower() or "main" in summary.lower()
            or "general" in tags or "main" in tags)


def generate_summary_and_tags(text, filename):
    llm = ChatOpenAI(temperature=0, model="gpt-3.5-turbo")  # 🔁 Updated model
    prompt = (
        f"Summarise the following document in 2–3 sentences. Then suggest up to 3 topical tags. "
        f"Document name: {filename}\n\n"
        f"Content:\n{text[:3000]}")
    response = llm.invoke(prompt).content.strip()

    summary_match = re.match(r"^(.*?)(?:Tags?:|\n|$)", response, re.DOTALL)
    tags_match = re.findall(
        r"#?(\b\w+\b)",
        response.split("Tags:")[-1]) if "Tags:" in response else []

    summary = summary_match.group(1).strip() if summary_match else ""
    tags = list(set(tag.lower() for tag in tags_match))

    if is_general_purpose(summary, tags, filename):
        tags.append("general")

    return summary, list(set(tags))


def process_blob(blob, container_client):
    filename = blob.name
    if not filename.endswith((".pdf", ".docx", ".pptx", ".txt")):
        return None

    blob_client = container_client.get_blob_client(filename)
    content = get_blob_content(blob_client)
    file_hash = calculate_file_hash(content)

    if existing_hashes.get(filename) == file_hash:
        print(f"✅ Skipped (no changes): {filename}")
        return None

    try:
        text = read_file(content, filename)
        summary, tags = generate_summary_and_tags(text, filename)

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=500,
            chunk_overlap=100)
        chunks = splitter.split_text(text)

        docs = [
            Document(page_content=chunk,
                     metadata={
                         "source": filename,
                         "summary": summary,
                         "tags": tags,
                     })
            for chunk in chunks
        ]

        print(f"✅ Processed: {filename}")
        return filename, file_hash, docs
    except Exception as e:
        print(f"❌ Failed to process {filename}: {e}")
        return None


def main():
    blob_service_client = BlobServiceClient.from_connection_string(
        AZURE_CONNECTION_STRING)
    container_client = blob_service_client.get_container_client(
        AZURE_CONTAINER_NAME)

    docs_with_metadata = []
    new_hashes = {}

    blobs = list(container_client.list_blobs())

    with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
        futures = [executor.submit(process_blob, blob, container_client)
                   for blob in blobs]
        for future in as_completed(futures):
            result = future.result()
            if result:
                filename, file_hash, docs = result
                docs_with_metadata.extend(docs)
                new_hashes[filename] = file_hash

    # ✅ SAFEGUARD: Only build FAISS if we have valid documents
    if not docs_with_metadata:
        print(
            "❌ No documents were successfully processed. FAISS index not built."
        )
        return

    print("✅ All documents processed. Now building FAISS index...")

    db = FAISS.from_documents(docs_with_metadata, OpenAIEmbeddings())
    db.save_local(INDEX_FILE)

    existing_hashes.update(new_hashes)
    with open(HASH_RECORD_FILE, "w") as f:
        json.dump(existing_hashes, f)

    print("✅ FAISS index saved.")


if __name__ == "__main__":
    main()
