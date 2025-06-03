import os
import json
from io import BytesIO
from flask import Flask, request, jsonify
from openai import OpenAI
from azure.storage.blob import BlobServiceClient
import fitz  # PyMuPDF
import docx  # python-docx
from pptx import Presentation  # python-pptx

# LangChain additions
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import OpenAIEmbeddings
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnableParallel, RunnablePassthrough
from langchain_openai import ChatOpenAI

app = Flask(__name__, static_folder="static")
client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))

# ── System prompt ──
SYSTEM_PERSONA = """
You are the “Greenwich Employability Embedding Assistant”, an employability adviser for the
Employability & Apprenticeship Directorate at the University of Greenwich (UK) with extensive experience.
Audience: academic staff (lecturers, module leaders, tutors) who want to embed
employability skills and resources into their curriculum.

Speaking style:
• British English / UK spelling only and academic-friendly tone.
• Professional and wholly relevant with specifics.
• Concise and clear.

Never mention internal implementation details (e.g., JSON, code).
"""

# ── Static fallback files ──
resources = {
    "cv": [
        {"name": "CV Template", "url": "/static/cv-template.pptx", "kind": "file"},
        {"name": "Employability Checklist", "url": "/static/checklist.pdf", "kind": "file"},
    ],
    "cover letter": [
        {"name": "Cover Letter Template", "url": "/static/cover-letter.docx", "kind": "file"}
    ]
}

# ── Load pathways.json ──
try:
    with open("pathways.json", "r", encoding="utf-8") as f:
        PATHWAYS = json.load(f)
except Exception as e:
    print("⚠️ Pathways JSON error →", e)
    PATHWAYS = []

def find_matching_pathways(query: str, limit: int = 5):
    q = query.lower()
    matches = []
    for p in PATHWAYS:
        if any(kw in q for kw in p.get("keywords", [])):
            matches.append({
                "name": p["title"],
                "url": p["url"],
                "info": p.get("description", ""),
                "kind": "pathway"
            })
        if len(matches) >= limit:
            break
    return matches

# ── Load FAISS vector store ──
try:
    db = FAISS.load_local("faiss_index", OpenAIEmbeddings())
except Exception as e:
    print("❌ FAISS index loading failed:", e)
    db = None

# ── Azure file reader ──
def download_blob_text(filename):
    try:
        conn_str = os.environ.get("AZURE_STORAGE_CONNECTION_STRING")
        blob_service = BlobServiceClient.from_connection_string(conn_str)
        container = blob_service.get_container_client("resources")
        blob = container.get_blob_client(filename)
        stream = BytesIO(blob.download_blob().readall())

        if filename.endswith(".pdf"):
            text = ""
            with fitz.open(stream=stream, filetype="pdf") as pdf:
                for page in pdf:
                    text += page.get_text()
            return text

        elif filename.endswith(".docx"):
            doc = docx.Document(stream)
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

        elif filename.endswith(".pptx"):
            prs = Presentation(stream)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text

        else:
            return "⚠️ Unsupported file type."

    except Exception as e:
        return f"❌ Error reading file: {e}"

# ── Routes ──
@app.route("/")
def index():
    return app.send_static_file("index.html")

@app.route("/ask", methods=["POST"])
def ask_gpt():
    data = request.get_json()
    user_message = data.get("message", "")

    if db is None:
        gpt_resp = client.chat.completions.create(model="gpt-3.5-turbo",
                                                  messages=[
                                                      {"role": "system", "content": SYSTEM_PERSONA},
                                                      {"role": "user", "content": user_message}
                                                  ])
        answer = gpt_resp.choices[0].message.content
    else:
        retriever = db.as_retriever(search_kwargs={"k": 5})

        prompt = ChatPromptTemplate.from_messages([
            ("system", SYSTEM_PERSONA),
            ("user", "{question}\n\nRelevant context:\n{context}")
        ])

        chain = (
            RunnableParallel({
                "context": retriever,
                "question": RunnablePassthrough()
            })
            | prompt
            | ChatOpenAI(model="gpt-3.5-turbo")
        )

        answer = chain.invoke(user_message).content

    matched = []
    for kw, files in resources.items():
        if kw in user_message.lower():
            matched.extend(files)

    matched.extend(find_matching_pathways(user_message))

    return jsonify({"reply": answer, "resources": matched})

@app.route("/read-file")
def read_file():
    filename = request.args.get("name")
    if not filename:
        return "⚠️ Please provide a filename using ?name=FILENAME"
    content = download_blob_text(filename)
    return f"<pre>{content}</pre>"

@app.route("/static/<path:filename>")
def static_files(filename):
    return app.send_static_file(filename)

@app.errorhandler(404)
def not_found(e):
    return app.send_static_file("index.html")

# ── Replit-specific local dev port ──
if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000)

